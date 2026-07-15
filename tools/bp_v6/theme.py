"""Editable PowerPoint primitives for the BP V6 visual system."""

from collections.abc import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt


SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Inches(7.5)
FONT_REGULAR = "Microsoft YaHei"
FONT_BOLD = "Microsoft YaHei"

NAVY = RGBColor(11, 24, 39)
NAVY_2 = RGBColor(20, 37, 58)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
PAPER = RGBColor(247, 249, 252)
WHITE = RGBColor(255, 255, 255)
CYAN = RGBColor(30, 153, 184)
CYAN_LIGHT = RGBColor(224, 246, 250)
AMBER = RGBColor(217, 119, 6)
AMBER_LIGHT = RGBColor(255, 247, 224)
RED = RGBColor(185, 28, 28)
RED_LIGHT = RGBColor(254, 236, 236)
GREEN = RGBColor(26, 127, 93)
GREEN_LIGHT = RGBColor(230, 247, 240)
LINE = RGBColor(214, 222, 232)


def new_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    return presentation


def set_background(slide, color=PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin / 2)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = FONT_BOLD if bold else FONT_REGULAR
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_lines(
    slide,
    lines: Iterable[tuple[str, bool, RGBColor]],
    x,
    y,
    w,
    h,
    *,
    size=14,
    bullet=False,
    spacing_after=5,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = Inches(0.04)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    for index, (text, bold, color) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = ""
        paragraph.space_after = Pt(spacing_after)
        if bullet:
            paragraph.text = "• "
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT_BOLD if bold else FONT_REGULAR
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, title: str, page: int):
    add_text(
        slide,
        f"{page:02d}",
        0.5,
        0.42,
        0.7,
        0.35,
        size=12,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    return add_text(
        slide,
        title,
        1.18,
        0.34,
        11.45,
        0.74,
        size=25,
        color=NAVY,
        bold=True,
    )


def add_card(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(0.8)
    return card


def add_badge(slide, text, x, y, w, *, fill=CYAN_LIGHT, color=CYAN):
    badge = add_card(slide, x, y, w, 0.34, fill=fill, line=fill)
    add_text(
        slide,
        text,
        x + 0.08,
        y + 0.055,
        w - 0.16,
        0.2,
        size=9,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    return badge


def add_rule(slide, x, y, w, *, color=LINE, width=1.0):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Pt(width),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_footer(slide, page: int):
    add_text(
        slide,
        "逸飞AI智眼系统｜创业大赛 BP V6",
        0.55,
        7.13,
        5.7,
        0.2,
        size=8,
        color=MUTED,
        margin=0,
    )
    add_text(
        slide,
        f"{page:02d}/16",
        10.75,
        7.13,
        2.0,
        0.2,
        size=8,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def add_source(slide, label: str, url: str, x: float, y: float, w: float):
    box = add_text(slide, label, x, y, w, 0.22, size=7, color=MUTED, margin=0)
    run = box.text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = url
    return box


def add_arrow(slide, x, y, w, h, *, fill=CYAN):
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill
    arrow.line.fill.background()
    return arrow
