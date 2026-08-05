"""Verify the final PPTX, PDF, workbook, and immutable source files."""

from __future__ import annotations

import argparse
import hashlib
import json
from pathlib import Path

import fitz
from openpyxl import load_workbook
from pptx import Presentation

from tools.bp_v6.content import SLIDES
from tools.bp_v6.excel_link import DEFAULT_PDF_NAME
from tools.bp_v6.theme import SLIDE_HEIGHT, SLIDE_WIDTH


FORBIDDEN_CLAIMS = (
    "100%自研",
    "完全自主知识产权",
    "国内领先",
    "已通过等保",
    "已服务多家单位",
    "已签约",
    "已落地",
)

REQUIRED_FACTS = (
    "公司成立以来收入\n0万元",
    "P95≤10秒",
    "试点验收目标，不是历史业绩",
    "融资需求\n500万元",
)


def sha256(path: str | Path) -> str:
    digest = hashlib.sha256()
    with Path(path).open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest().upper()


def _pptx_text(presentation: Presentation) -> str:
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def _verify_sources(baseline: dict) -> bool:
    for key in ("source_pdf", "source_workbook"):
        item = baseline[key]
        path = Path(item["path"])
        if not path.is_file():
            raise FileNotFoundError(path)
        if sha256(path) != item["sha256"].upper():
            raise ValueError(f"immutable source changed: {path}")
    return True


def _verify_pptx(path: Path) -> dict[str, object]:
    presentation = Presentation(path)
    if len(presentation.slides) != 16:
        raise ValueError(f"PPTX must contain 16 slides, found {len(presentation.slides)}")
    if presentation.slide_width != SLIDE_WIDTH or presentation.slide_height != SLIDE_HEIGHT:
        raise ValueError("PPTX slide size is not the approved 16:9 canvas")
    text = _pptx_text(presentation)
    for spec in SLIDES:
        if spec.title not in text:
            raise ValueError(f"PPTX is missing slide title: {spec.title}")
    for fact in REQUIRED_FACTS:
        if fact not in text:
            raise ValueError(f"PPTX is missing locked fact: {fact}")
    for claim in FORBIDDEN_CLAIMS:
        if claim in text:
            raise ValueError(f"PPTX contains forbidden claim: {claim}")
    return {
        "path": str(path.resolve()),
        "sha256": sha256(path),
        "slides": len(presentation.slides),
        "width_emu": presentation.slide_width,
        "height_emu": presentation.slide_height,
    }


def _verify_pdf(path: Path) -> dict[str, object]:
    document = fitz.open(path)
    try:
        if document.page_count != 16:
            raise ValueError(f"PDF must contain 16 pages, found {document.page_count}")
        first_page = document[0]
        ratio = first_page.rect.width / first_page.rect.height
        if abs(ratio - 16 / 9) > 0.03:
            raise ValueError(f"PDF page ratio is not 16:9: {ratio:.4f}")
        extracted_chars = sum(len(page.get_text("text").strip()) for page in document)
        if extracted_chars < 64:
            raise ValueError("PDF text extraction is unexpectedly empty")
    finally:
        document.close()
    return {
        "path": str(path.resolve()),
        "sha256": sha256(path),
        "pages": 16,
        "aspect_ratio": ratio,
        "extracted_text_characters": extracted_chars,
    }


def _verify_workbook(path: Path, baseline: dict, pdf_path: Path) -> dict[str, object]:
    workbook = load_workbook(path, data_only=False, keep_links=True)
    try:
        expected_count = baseline["source_workbook"]["sheet_count"]
        expected_first = baseline["source_workbook"]["first_sheet"]
        cell_ref = baseline["source_workbook"]["pdf_cell"]
        if len(workbook.worksheets) != expected_count:
            raise ValueError("final workbook sheet count changed")
        if workbook.worksheets[0].title != expected_first:
            raise ValueError("final workbook first sheet changed")
        cell = workbook.worksheets[0][cell_ref]
        target = cell.hyperlink.target if cell.hyperlink else None
        if cell.value != DEFAULT_PDF_NAME or target != DEFAULT_PDF_NAME:
            raise ValueError("final workbook PDF display value or hyperlink is incorrect")
    finally:
        workbook.close()
    linked_pdf = path.parent / DEFAULT_PDF_NAME
    if linked_pdf.resolve() != pdf_path.resolve() or not linked_pdf.is_file():
        raise ValueError("final workbook relative PDF target is missing")
    return {
        "path": str(path.resolve()),
        "sha256": sha256(path),
        "sheets": expected_count,
        "first_sheet": expected_first,
        "pdf_cell": cell_ref,
        "pdf_link": DEFAULT_PDF_NAME,
    }


def verify_artifacts(
    pptx_path: str | Path,
    pdf_path: str | Path,
    workbook_path: str | Path,
    baseline_path: str | Path,
) -> dict[str, object]:
    """Run the complete deterministic artifact contract."""

    pptx = Path(pptx_path)
    pdf = Path(pdf_path)
    workbook = Path(workbook_path)
    for path in (pptx, pdf, workbook):
        if not path.is_file():
            raise FileNotFoundError(path)
    baseline = json.loads(Path(baseline_path).read_text(encoding="utf-8-sig"))
    report = {
        "status": "passed",
        "sources_unchanged": _verify_sources(baseline),
        "pptx": _verify_pptx(pptx),
        "pdf": _verify_pdf(pdf),
        "workbook": _verify_workbook(workbook, baseline, pdf),
    }
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--pdf", type=Path, required=True)
    parser.add_argument("--workbook", type=Path, required=True)
    parser.add_argument("--baseline", type=Path, required=True)
    parser.add_argument("--report", type=Path)
    args = parser.parse_args()
    report = verify_artifacts(args.pptx, args.pdf, args.workbook, args.baseline)
    payload = json.dumps(report, ensure_ascii=False, indent=2)
    if args.report:
        args.report.parent.mkdir(parents=True, exist_ok=True)
        args.report.write_text(payload, encoding="utf-8")
    print(payload)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
