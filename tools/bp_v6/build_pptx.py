"""Build the editable 16-slide yFeiEye BP V6 presentation."""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Callable

from openpyxl import load_workbook
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from tools.bp_v6.assets import EXPECTED_ASSETS, validate_assets
from tools.bp_v6.content import (
    ARCHITECTURE_NODES,
    COMPARISON_ROWS,
    COMPANY,
    COMPLIANCE_CONTROLS,
    FORECAST,
    FOUNDER,
    FUNDING,
    MARKET,
    MILESTONES,
    PILOT_METRICS,
    PRICING_ROWS,
    PROCUREMENT_EXAMPLES,
    SECURITY_CONTROLS,
    SLIDES,
    SOURCES,
)
from tools.bp_v6.theme import (
    AMBER,
    AMBER_LIGHT,
    CYAN,
    CYAN_LIGHT,
    GREEN,
    GREEN_LIGHT,
    INK,
    LINE,
    MUTED,
    NAVY,
    NAVY_2,
    PAPER,
    RED,
    RED_LIGHT,
    WHITE,
    add_arrow,
    add_badge,
    add_card,
    add_footer,
    add_rule,
    add_source,
    add_text,
    add_title,
    new_presentation,
    set_background,
)


LIGHT_BLUE = RGBColor(234, 242, 249)
LIGHT_NAVY = RGBColor(222, 229, 238)
MID_BLUE = RGBColor(74, 115, 145)


def _shape(slide, shape_type, x, y, w, h, fill, *, line=None, width=0.8):
    item = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    item.fill.solid()
    item.fill.fore_color.rgb = fill
    if line is None:
        item.line.fill.background()
    else:
        item.line.color.rgb = line
        item.line.width = Pt(width)
    return item


def _circle_label(slide, value, x, y, diameter=0.48, *, fill=CYAN, color=WHITE):
    _shape(slide, MSO_SHAPE.OVAL, x, y, diameter, diameter, fill)
    add_text(
        slide,
        value,
        x,
        y + 0.07,
        diameter,
        diameter - 0.1,
        size=11,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def _section_slide(prs, spec, *, background=PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, background)
    add_title(slide, spec.title, spec.number)
    add_footer(slide, spec.number)
    return slide


def _card_copy(
    slide,
    x,
    y,
    w,
    h,
    title,
    body,
    *,
    accent=CYAN,
    fill=WHITE,
    title_size=14,
    body_size=10,
):
    add_card(slide, x, y, w, h, fill=fill)
    _shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.06, h, accent)
    add_text(slide, title, x + 0.24, y + 0.2, w - 0.42, 0.35, size=title_size, color=NAVY, bold=True)
    add_text(slide, body, x + 0.24, y + 0.68, w - 0.42, h - 0.84, size=body_size, color=MUTED, line_spacing=1.05)


def _metric_card(slide, x, y, w, label, value, note, *, accent=CYAN):
    add_card(slide, x, y, w, 1.48, fill=WHITE)
    add_text(slide, label, x + 0.2, y + 0.16, w - 0.4, 0.24, size=9, color=MUTED, bold=True)
    add_text(slide, value, x + 0.2, y + 0.49, w - 0.4, 0.44, size=22, color=accent, bold=True)
    add_text(slide, note, x + 0.2, y + 1.04, w - 0.4, 0.28, size=8, color=MUTED)


def _sources(slide, spec, y=6.73):
    keys = spec.source_keys
    if not keys:
        return
    count = len(keys)
    columns = 4 if count > 4 else count
    width = 12.15 / max(columns, 1)
    for index, key in enumerate(keys):
        source = SOURCES[key]
        row = index // columns
        col = index % columns
        label = source.label
        if len(label) > 24:
            label = label[:23] + "…"
        add_source(
            slide,
            f"来源：{label}",
            source.url,
            0.58 + col * width,
            y + row * 0.18,
            width - 0.12,
        )


def _picture_crop(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    frame_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio > frame_ratio:
        visible = frame_ratio / image_ratio
        crop = (1 - visible) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < frame_ratio:
        visible = image_ratio / frame_ratio
        crop = (1 - visible) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def _cover(prs, spec, assets, contact):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, 7.5, CYAN)
    _shape(slide, MSO_SHAPE.RECTANGLE, 8.8, 0, 4.53, 7.5, NAVY_2)
    for i, width in enumerate((3.5, 2.85, 2.25, 1.65)):
        _shape(slide, MSO_SHAPE.RECTANGLE, 9.25, 1.1 + i * 1.06, width, 0.05, MID_BLUE)
        _circle_label(slide, f"0{i + 1}", 9.0, 0.89 + i * 1.06, 0.45, fill=CYAN if i == 0 else MID_BLUE)
    add_badge(slide, "AI VIDEO · CLOSED LOOP", 0.7, 0.6, 2.35, fill=CYAN, color=WHITE)
    add_text(slide, spec.title, 0.7, 1.45, 7.65, 1.1, size=38, color=WHITE, bold=True)
    add_text(slide, spec.body[0], 0.7, 2.72, 6.0, 0.52, size=23, color=CYAN, bold=True)
    add_text(slide, spec.body[1], 0.7, 3.52, 7.15, 0.86, size=16, color=LIGHT_NAVY, line_spacing=1.12)
    add_rule(slide, 0.7, 4.65, 7.2, color=MID_BLUE)
    tags = spec.body[2].split("｜")
    for index, tag in enumerate(tags):
        add_text(slide, tag, 0.72 + index * 1.86, 4.93, 1.65, 0.3, size=10, color=WHITE, bold=True)
        if index < len(tags) - 1:
            _shape(slide, MSO_SHAPE.OVAL, 2.25 + index * 1.86, 5.02, 0.07, 0.07, CYAN)
    add_text(slide, COMPANY["name"], 0.7, 6.55, 5.4, 0.3, size=10, color=LIGHT_NAVY)
    add_text(slide, "创业大赛商业计划书 · V6", 9.25, 5.72, 3.1, 0.36, size=13, color=WHITE, bold=True)
    add_text(slide, "2026.07", 9.25, 6.18, 2.0, 0.28, size=10, color=CYAN, bold=True)
    add_text(slide, "01/16", 10.75, 7.13, 2.0, 0.2, size=8, color=LIGHT_NAVY, align=PP_ALIGN.RIGHT, margin=0)


def _pain(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    cards = (
        ("01", "告警孤岛", "识别结果分散，难形成统一事件视图", RED, RED_LIGHT),
        ("02", "人工盯屏", "画面多、班次长，异常发现依赖经验", AMBER, AMBER_LIGHT),
        ("03", "重复误报", "同类告警反复出现，消耗复核与处置精力", AMBER, AMBER_LIGHT),
        ("04", "处置断点", "派发、回执、复查与归档缺少连续证据", RED, RED_LIGHT),
    )
    for index, (number, title, body, accent, fill) in enumerate(cards):
        x = 0.65 + (index % 2) * 6.12
        y = 1.45 + (index // 2) * 1.65
        add_card(slide, x, y, 5.75, 1.32, fill=WHITE)
        _circle_label(slide, number, x + 0.24, y + 0.22, 0.5, fill=accent)
        add_text(slide, title, x + 0.93, y + 0.18, 1.75, 0.34, size=16, color=NAVY, bold=True)
        add_text(slide, body, x + 0.93, y + 0.64, 4.35, 0.4, size=10, color=MUTED)
    add_card(slide, 0.65, 4.93, 12.05, 1.3, fill=NAVY, line=NAVY)
    add_text(slide, "采购信号 ≠ 客户背书", 0.96, 5.2, 2.55, 0.34, size=15, color=CYAN, bold=True)
    add_text(slide, spec.body[-1], 3.5, 5.12, 8.75, 0.62, size=11, color=WHITE, line_spacing=1.08)
    add_text(slide, "公开案例预算：90万 / 209.6万 / 332万 / 2836.9万", 3.5, 5.82, 8.5, 0.24, size=8, color=LIGHT_NAVY)
    _sources(slide, spec, y=6.57)


def _workflow(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    labels = spec.body
    if spec.number == 3:
        accents = (MID_BLUE, MID_BLUE, CYAN, CYAN, AMBER, GREEN, NAVY)
        y = 2.05
        for index, (label, accent) in enumerate(zip(labels, accents, strict=True)):
            x = 0.55 + index * 1.82
            _circle_label(slide, f"{index + 1:02d}", x + 0.45, y - 0.57, 0.48, fill=accent)
            add_card(slide, x, y, 1.48, 1.15, fill=WHITE)
            add_text(slide, label, x + 0.12, y + 0.34, 1.24, 0.4, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            if index < 6:
                add_arrow(slide, x + 1.47, y + 0.42, 0.34, 0.25, fill=LINE)
        add_card(slide, 0.65, 4.24, 12.05, 1.25, fill=CYAN_LIGHT, line=CYAN_LIGHT)
        add_text(slide, "产品价值不是多一个告警窗口", 0.92, 4.5, 3.25, 0.35, size=15, color=CYAN, bold=True)
        add_text(slide, "而是让一次AI命中拥有复核责任、处置时限、结果证据和完整审计链。", 4.2, 4.45, 7.9, 0.52, size=14, color=NAVY, bold=True)
        add_text(slide, "AI只提供辅助判断；授权人员确认前，不触发外部处置动作。", 4.2, 5.06, 7.6, 0.26, size=9, color=MUTED)
    else:
        steps = labels[:5]
        for index, line in enumerate(steps):
            title, _, body = line.partition("：")
            x = 0.58 + index * 2.52
            _circle_label(slide, f"{index + 1:02d}", x + 0.72, 1.55, 0.5, fill=CYAN if index < 3 else GREEN)
            add_card(slide, x, 2.18, 2.12, 2.35, fill=WHITE)
            add_text(slide, title, x + 0.18, 2.48, 1.76, 0.36, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, body, x + 0.22, 3.12, 1.68, 0.86, size=9, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.08)
            if index < 4:
                add_arrow(slide, x + 2.03, 3.02, 0.47, 0.3, fill=LINE)
        add_card(slide, 0.65, 5.05, 12.05, 0.9, fill=AMBER_LIGHT, line=AMBER_LIGHT)
        add_text(slide, "当前边界", 0.95, 5.31, 1.15, 0.25, size=11, color=AMBER, bold=True)
        add_text(slide, labels[-1], 2.12, 5.25, 4.4, 0.34, size=14, color=NAVY, bold=True)
        add_text(slide, "先用技术交流建立场景共识，再以付费试点形成可验证证据。", 6.55, 5.3, 5.55, 0.28, size=10, color=MUTED)


def _stage(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    values = (
        ("成立日期", "2025-12-09", CYAN),
        ("公司成立以来收入", "0万元", AMBER),
        ("正式客户 / 合同 / 试点", "0 / 0 / 0", RED),
        ("核心团队", "1人", CYAN),
    )
    for index, (label, value, accent) in enumerate(values):
        _metric_card(slide, 0.65 + index * 3.05, 1.42, 2.75, label, value, "已核实事实", accent=accent)
    add_text(slide, "证据成熟度", 0.67, 3.28, 2.0, 0.32, size=14, color=NAVY, bold=True)
    stages = (
        ("01", "平台可演示", "设备、任务、复核链路", GREEN),
        ("02", "交付可复用", "报价、清单、指标、模板", CYAN),
        ("03", "付费试点", "8路 · 30天 · 共同验收", AMBER),
        ("04", "规模复制", "单所扩容与多点复制", MID_BLUE),
    )
    for index, (number, title, note, accent) in enumerate(stages):
        x = 0.65 + index * 3.04
        add_card(slide, x, 3.85, 2.72, 1.48, fill=WHITE)
        _circle_label(slide, number, x + 0.18, 4.13, 0.43, fill=accent)
        add_text(slide, title, x + 0.78, 4.04, 1.65, 0.31, size=13, color=NAVY, bold=True)
        add_text(slide, note, x + 0.78, 4.53, 1.64, 0.38, size=9, color=MUTED)
        if index < 3:
            add_arrow(slide, x + 2.66, 4.39, 0.4, 0.26, fill=LINE)
    add_card(slide, 0.65, 5.68, 12.05, 0.72, fill=NAVY, line=NAVY)
    add_text(slide, "下一证据", 0.92, 5.91, 1.25, 0.24, size=11, color=CYAN, bold=True)
    add_text(slide, "用8路、30天付费试点验证闭环价值与交付方法", 2.15, 5.84, 8.0, 0.34, size=15, color=WHITE, bold=True)


def _screenshots(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    items = (
        ("设备接入", "设备-01｜演示通道-01", EXPECTED_ASSETS[1]),
        ("算法任务", "司法场景任务-01｜运行中", EXPECTED_ASSETS[0]),
        ("线索复核", "确认前不会执行任何动作", EXPECTED_ASSETS[2]),
    )
    for index, (title, caption, filename) in enumerate(items):
        x = 0.56 + index * 4.24
        add_card(slide, x, 1.37, 4.02, 4.48, fill=WHITE)
        add_text(slide, title, x + 0.2, 1.58, 2.15, 0.32, size=14, color=NAVY, bold=True)
        add_badge(slide, "真实界面 · 已脱敏", x + 2.4, 1.55, 1.34, fill=CYAN_LIGHT, color=CYAN)
        picture = _picture_crop(slide, assets / filename, x + 0.16, 2.06, 3.7, 2.08)
        if filename == EXPECTED_ASSETS[2]:
            picture.crop_top = 0
            picture.crop_bottom = 0.46
        add_text(slide, caption, x + 0.2, 4.36, 3.62, 0.32, size=10, color=NAVY, bold=True)
        add_text(slide, spec.body[index], x + 0.2, 4.84, 3.58, 0.66, size=8, color=MUTED, line_spacing=1.05)
    add_card(slide, 0.65, 6.08, 12.05, 0.55, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "说明", 0.9, 6.25, 0.55, 0.2, size=9, color=AMBER, bold=True)
    add_text(slide, spec.body[-1], 1.5, 6.2, 10.7, 0.27, size=8, color=MUTED)


def _architecture(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    add_card(slide, 0.6, 1.35, 12.15, 4.94, fill=WHITE, line=CYAN)
    add_badge(slide, "客户内网 / 专有云安全边界", 0.9, 1.56, 2.15, fill=CYAN, color=WHITE)
    main_nodes = ARCHITECTURE_NODES[:4]
    fills = (LIGHT_BLUE, CYAN_LIGHT, CYAN_LIGHT, GREEN_LIGHT)
    for index, (label, fill) in enumerate(zip(main_nodes, fills, strict=True)):
        x = 0.95 + index * 2.8
        add_card(slide, x, 2.18, 2.33, 1.2, fill=fill, line=fill)
        add_text(slide, label, x + 0.15, 2.46, 2.03, 0.62, size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if index < 3:
            add_arrow(slide, x + 2.32, 2.63, 0.48, 0.29, fill=LINE)
    add_arrow(slide, 5.99, 3.57, 0.38, 0.62, fill=CYAN)
    add_card(slide, 3.25, 4.22, 6.83, 1.16, fill=NAVY, line=NAVY)
    add_text(slide, ARCHITECTURE_NODES[5], 3.55, 4.5, 6.23, 0.57, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_badge(slide, "告警 → 事件", 0.98, 4.39, 1.58, fill=AMBER_LIGHT, color=AMBER)
    add_text(slide, "AI命中只进入复核队列", 0.95, 4.9, 1.78, 0.4, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_badge(slide, "监管终端", 10.75, 4.39, 1.38, fill=GREEN_LIGHT, color=GREEN)
    add_text(slide, "驾驶舱 / 待办 / 审计", 10.58, 4.9, 1.74, 0.4, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    add_rule(slide, 0.95, 5.69, 11.48, color=LINE)
    for index, control in enumerate(SECURITY_CONTROLS):
        add_text(slide, control, 0.95 + index * 2.29, 5.85, 2.0, 0.25, size=8, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "贯穿：身份权限 · 审计 · 导出审批 · 保留删除 · 备份恢复", 2.72, 6.47, 7.9, 0.24, size=9, color=CYAN, bold=True, align=PP_ALIGN.CENTER)


def _comparison(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    x0, y0 = 0.62, 1.42
    widths = (2.25, 2.75, 2.75, 4.1)
    headers = ("比较维度", "传统视频 / SI", "单点算法厂商", "逸飞AI智眼系统")
    x = x0
    for index, (header, width) in enumerate(zip(headers, widths, strict=True)):
        fill = NAVY if index < 3 else CYAN
        _shape(slide, MSO_SHAPE.RECTANGLE, x, y0, width, 0.66, fill)
        add_text(slide, header, x + 0.08, y0 + 0.19, width - 0.16, 0.25, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += width
    row_h = 0.78
    for row_index, row in enumerate(COMPARISON_ROWS):
        x = x0
        y = y0 + 0.66 + row_index * row_h
        for col_index, (value, width) in enumerate(zip(row, widths, strict=True)):
            fill = CYAN_LIGHT if col_index == 3 else (WHITE if row_index % 2 == 0 else PAPER)
            _shape(slide, MSO_SHAPE.RECTANGLE, x, y, width, row_h, fill, line=LINE, width=0.4)
            add_text(slide, value, x + 0.12, y + 0.16, width - 0.24, row_h - 0.22, size=8 if col_index else 9, color=NAVY if col_index in (0, 3) else MUTED, bold=col_index in (0, 3), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
            x += width
    add_card(slide, 0.62, 6.12, 11.85, 0.48, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, spec.body[-1], 0.85, 6.26, 11.3, 0.2, size=8, color=MUTED, align=PP_ALIGN.CENTER)


def _compliance(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    for index, (title, body) in enumerate(COMPLIANCE_CONTROLS):
        x = 0.65 + (index % 3) * 4.05
        y = 1.38 + (index // 3) * 2.12
        accent = CYAN if index < 3 else GREEN
        _card_copy(slide, x, y, 3.75, 1.72, title, body, accent=accent, title_size=13, body_size=9)
    add_card(slide, 0.65, 5.84, 12.05, 0.58, fill=RED_LIGHT, line=RED_LIGHT)
    add_text(slide, "边界声明", 0.88, 6.01, 1.0, 0.22, size=9, color=RED, bold=True)
    add_text(slide, spec.body[-1], 1.85, 5.97, 10.3, 0.28, size=8, color=MUTED)
    _sources(slide, spec, y=6.58)


def _pilot(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    add_card(slide, 0.65, 1.4, 12.05, 0.77, fill=NAVY, line=NAVY)
    scopes = (("8路", "视频接入"), ("30天", "验证周期"), ("2-3类", "高价值规则"), ("共同", "指标验收"))
    for index, (value, label) in enumerate(scopes):
        x = 0.95 + index * 2.93
        add_text(slide, value, x, 1.57, 1.0, 0.31, size=17, color=CYAN, bold=True)
        add_text(slide, label, x + 1.02, 1.64, 1.25, 0.22, size=9, color=WHITE, bold=True)
    accents = (CYAN, CYAN, GREEN, AMBER)
    for index, ((label, value, note), accent) in enumerate(zip(PILOT_METRICS, accents, strict=True)):
        _metric_card(slide, 0.65 + index * 3.05, 2.58, 2.75, label, value, note, accent=accent)
    add_text(slide, "30天验证节奏", 0.67, 4.55, 2.0, 0.3, size=13, color=NAVY, bold=True)
    phases = (("D1-3", "基线与规则"), ("D4-14", "运行与周报"), ("D15-27", "误报复盘"), ("D28-30", "复查与验收"))
    for index, (time, label) in enumerate(phases):
        x = 0.66 + index * 3.02
        _shape(slide, MSO_SHAPE.RECTANGLE, x, 5.08, 2.72, 0.1, CYAN if index < 3 else GREEN)
        add_text(slide, time, x, 5.35, 0.9, 0.27, size=10, color=CYAN if index < 3 else GREEN, bold=True)
        add_text(slide, label, x + 0.88, 5.35, 1.75, 0.27, size=10, color=NAVY, bold=True)
    add_badge(slide, "试点验收目标，不是历史业绩", 4.72, 6.17, 3.9, fill=AMBER_LIGHT, color=AMBER)


def _market(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    funnel = (
        ("TAM背景", ">9600亿元", "2024安防行业总产值", LIGHT_NAVY, 5.5),
        ("SAM测算", "7.9-16.4亿元", "1051所 × 75.4-156.2万元", CYAN_LIGHT, 4.65),
        ("SOM规划", "2350万元", "2027-2029累计目标", GREEN_LIGHT, 3.8),
    )
    for index, (label, value, note, fill, width) in enumerate(funnel):
        x = 0.7 + (5.55 - width) / 2
        y = 1.55 + index * 1.28
        _shape(slide, MSO_SHAPE.TRAPEZOID, x, y, width, 1.03, fill, line=LINE)
        add_text(slide, label, x + 0.35, y + 0.18, 1.05, 0.25, size=9, color=MUTED, bold=True)
        add_text(slide, value, x + 1.47, y + 0.13, width - 1.8, 0.34, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, note, x + 0.45, y + 0.61, width - 0.9, 0.19, size=7, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "公开采购预算信号", 6.57, 1.48, 2.5, 0.32, size=14, color=NAVY, bold=True)
    for index, (name, amount, note) in enumerate(PROCUREMENT_EXAMPLES):
        x = 6.55 + (index % 2) * 3.05
        y = 1.96 + (index // 2) * 1.56
        add_card(slide, x, y, 2.78, 1.27, fill=WHITE)
        amount_size = 17 if len(amount) <= 6 else 14
        add_text(slide, amount, x + 0.18, y + 0.17, 2.35, 0.34, size=amount_size, color=CYAN, bold=True)
        add_text(slide, name, x + 0.18, y + 0.59, 2.35, 0.27, size=9, color=NAVY, bold=True)
        add_text(slide, note, x + 0.18, y + 0.91, 2.2, 0.18, size=7, color=MUTED)
    add_card(slide, 0.65, 5.3, 12.05, 0.88, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "测算口径", 0.88, 5.54, 0.95, 0.24, size=9, color=AMBER, bold=True)
    add_text(slide, spec.body[-1], 1.8, 5.45, 10.25, 0.42, size=8, color=MUTED)
    _sources(slide, spec, y=6.39)


def _pricing(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    for index, (name, price, note) in enumerate(PRICING_ROWS):
        x = 0.62 + (index % 3) * 4.08
        y = 1.36 + (index // 3) * 2.2
        fill = CYAN_LIGHT if index == 0 else WHITE
        accent = CYAN if index < 4 else GREEN
        add_card(slide, x, y, 3.78, 1.76, fill=fill)
        add_text(slide, name, x + 0.22, y + 0.2, 2.2, 0.31, size=13, color=NAVY, bold=True)
        add_text(slide, f"{price:g}", x + 0.2, y + 0.61, 1.7, 0.52, size=26, color=accent, bold=True)
        add_text(slide, "万元" + ("/年" if index == 4 else "/类" if index == 5 else ""), x + 1.73, y + 0.83, 0.85, 0.23, size=8, color=MUTED, bold=True)
        add_text(slide, note, x + 0.22, y + 1.28, 3.22, 0.25, size=8, color=MUTED)
    add_card(slide, 0.65, 5.93, 12.05, 0.5, fill=NAVY, line=NAVY)
    add_text(slide, spec.body[-1], 0.9, 6.08, 11.5, 0.2, size=8, color=WHITE, align=PP_ALIGN.CENTER)


def _founder(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    add_card(slide, 0.65, 1.38, 3.4, 4.92, fill=NAVY, line=NAVY)
    _shape(slide, MSO_SHAPE.OVAL, 1.66, 1.83, 1.38, 1.38, CYAN)
    add_text(slide, "刘飞", 1.66, 2.24, 1.38, 0.45, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    education = FOUNDER["education"].replace("｜计算机科学与技术", "\n计算机科学与技术")
    add_text(slide, education, 0.98, 3.55, 2.75, 0.72, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rule(slide, 1.05, 4.48, 2.6, color=MID_BLUE)
    add_text(slide, "7年", 1.0, 4.8, 0.8, 0.34, size=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "5年", 1.95, 4.8, 0.8, 0.34, size=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "13年", 2.9, 4.8, 0.8, 0.34, size=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "软件研发", 1.0, 5.24, 0.8, 0.2, size=8, color=LIGHT_NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, "项目经理", 1.95, 5.24, 0.8, 0.2, size=8, color=LIGHT_NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, "从业经验", 2.9, 5.24, 0.8, 0.2, size=8, color=LIGHT_NAVY, align=PP_ALIGN.CENTER)
    add_badge(slide, "当前核心团队 1人", 1.17, 5.7, 2.35, fill=CYAN, color=WHITE)
    add_text(slide, "当前分工", 4.42, 1.43, 2.0, 0.33, size=14, color=NAVY, bold=True)
    for index, role in enumerate(FOUNDER["responsibilities"]):
        x = 4.4 + (index % 2) * 2.98
        y = 1.94 + (index // 2) * 1.15
        _card_copy(slide, x, y, 2.7, 0.92, f"0{index + 1}", role, accent=CYAN, title_size=9, body_size=9)
    add_text(slide, "行业资源与边界", 10.55, 1.43, 2.0, 0.33, size=14, color=NAVY, bold=True)
    add_card(slide, 10.43, 1.94, 2.28, 2.07, fill=CYAN_LIGHT, line=CYAN_LIGHT)
    add_text(slide, "多次行业峰会\n司法安防技术交流\n多年检测与视频经验", 10.68, 2.27, 1.78, 1.13, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.18)
    add_text(slide, "单位名称按保密要求不披露；技术交流不等于客户或订单。", 10.65, 3.49, 1.84, 0.36, size=7, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, 4.4, 4.52, 8.3, 1.53, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "融资后补位", 4.7, 4.8, 1.25, 0.31, size=13, color=AMBER, bold=True)
    hires = ("算法 / 数据", "司法方案 / 交付", "销售 / 渠道", "安全 / 合规")
    for index, role in enumerate(hires):
        add_badge(slide, role, 6.0 + index * 1.56, 4.74, 1.42, fill=WHITE, color=NAVY)
    add_text(slide, "创始人现阶段对产品、架构与交付证据负责；扩张以前先补齐关键岗位。", 4.72, 5.42, 7.4, 0.25, size=8, color=MUTED)


def _timeline(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    tracks = (("Q3 2026", MILESTONES["Q3 2026"], CYAN), ("Q4 2026", MILESTONES["Q4 2026"], GREEN))
    for col, (quarter, items, accent) in enumerate(tracks):
        x = 0.65 + col * 6.1
        add_card(slide, x, 1.42, 5.78, 4.48, fill=WHITE)
        add_badge(slide, quarter, x + 0.28, 1.7, 1.38, fill=accent, color=WHITE)
        for index, item in enumerate(items):
            y = 2.38 + index * 0.81
            _circle_label(slide, f"{index + 1:02d}", x + 0.32, y, 0.38, fill=accent)
            add_text(slide, item, x + 0.88, y - 0.01, 4.45, 0.5, size=9, color=NAVY, bold=index == 0, line_spacing=1.04)
    add_card(slide, 0.65, 6.12, 12.05, 0.49, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "以上均为计划和目标，不作为当前已完成事项。", 0.9, 6.27, 11.55, 0.2, size=8, color=AMBER, bold=True, align=PP_ALIGN.CENTER)


def _finance(prs, spec, assets, contact):
    slide = _section_slide(prs, spec)
    add_text(slide, "三年收入规划（万元）", 0.65, 1.36, 3.1, 0.32, size=14, color=NAVY, bold=True)
    max_value = max(item["target"] for item in FORECAST.values())
    for index, (year, values) in enumerate(FORECAST.items()):
        y = 2.02 + index * 1.15
        add_text(slide, str(year), 0.68, y + 0.12, 0.63, 0.25, size=10, color=MUTED, bold=True)
        bar_w = 4.7 * values["target"] / max_value
        _shape(slide, MSO_SHAPE.RECTANGLE, 1.42, y, 4.7, 0.57, LIGHT_NAVY)
        _shape(slide, MSO_SHAPE.RECTANGLE, 1.42, y, max(bar_w, 0.62), 0.57, CYAN if index < 2 else GREEN)
        add_text(slide, f"{values['target']} 目标", 1.62, y + 0.13, 1.2, 0.23, size=9, color=WHITE, bold=True)
        add_text(slide, f"测算 {values['calculated']}｜{values['mix']}", 1.43, y + 0.7, 4.8, 0.24, size=7, color=MUTED)
    add_text(slide, "融资需求\n500万元", 7.15, 1.43, 2.15, 0.92, size=23, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "资金用途", 9.67, 1.5, 1.5, 0.3, size=13, color=NAVY, bold=True)
    colors = (CYAN, MID_BLUE, GREEN, AMBER, RED)
    total_w = 5.55
    x = 6.75
    for item, color in zip(FUNDING, colors, strict=True):
        width = total_w * item["ratio"] / 100
        _shape(slide, MSO_SHAPE.RECTANGLE, x, 2.64, width, 0.52, color)
        if item["ratio"] >= 15:
            add_text(slide, f"{item['ratio']}%", x, 2.78, width, 0.18, size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
        x += width
    for index, (item, color) in enumerate(zip(FUNDING, colors, strict=True)):
        y = 3.48 + index * 0.53
        _shape(slide, MSO_SHAPE.OVAL, 6.78, y + 0.02, 0.15, 0.15, color)
        add_text(slide, item["name"], 7.05, y, 3.55, 0.22, size=8, color=NAVY, bold=True)
        add_text(slide, f"{item['ratio']}%｜{item['amount']}万元", 10.75, y, 1.35, 0.22, size=8, color=MUTED, align=PP_ALIGN.RIGHT)
    add_card(slide, 0.65, 6.28, 12.05, 0.36, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "规划目标，不是订单储备或历史业绩；基础预测不计年度治理续费。", 0.9, 6.38, 11.5, 0.16, size=7, color=AMBER, bold=True, align=PP_ALIGN.CENTER)


def _cta(prs, spec, assets, contact):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, 7.5, CYAN)
    add_badge(slide, "NEXT PROOF", 0.7, 0.62, 1.42, fill=CYAN, color=WHITE)
    add_text(slide, spec.title, 0.7, 1.35, 11.3, 0.86, size=32, color=WHITE, bold=True)
    add_text(slide, "不是讲更多故事，而是共同产生第一份可核验业务证据。", 0.72, 2.45, 9.3, 0.4, size=14, color=LIGHT_NAVY)
    asks = (
        ("01", "付费试点", "8路 · 30天\n形成共同验收证据"),
        ("02", "集成协同", "接口适配 · 现场部署\n联合交付与客户验收"),
        ("03", "能力资源", "算力 · 数据治理\n产业资源与融资支持"),
    )
    for index, (number, title, body) in enumerate(asks):
        x = 0.72 + index * 4.03
        add_card(slide, x, 3.18, 3.68, 1.78, fill=NAVY_2, line=MID_BLUE)
        _circle_label(slide, number, x + 0.25, 3.48, 0.5, fill=CYAN)
        add_text(slide, title, x + 0.92, 3.38, 2.15, 0.34, size=16, color=WHITE, bold=True)
        add_text(slide, body, x + 0.92, 3.92, 2.35, 0.65, size=10, color=LIGHT_NAVY, line_spacing=1.1)
    add_rule(slide, 0.72, 5.51, 11.92, color=MID_BLUE)
    add_text(slide, spec.body[-1], 0.72, 5.84, 8.7, 0.45, size=18, color=CYAN, bold=True)
    name = contact.get("name", "")
    phone = contact.get("phone", "")
    email = contact.get("email", "")
    contact_text = "｜".join(value for value in (name, phone, email) if value)
    if contact_text:
        add_text(slide, f"联系：{contact_text}", 0.72, 6.55, 9.6, 0.28, size=10, color=WHITE, bold=True)
    add_text(slide, COMPANY["name"], 9.45, 6.55, 3.18, 0.28, size=9, color=LIGHT_NAVY, align=PP_ALIGN.RIGHT)
    add_text(slide, "16/16", 10.75, 7.13, 2.0, 0.2, size=8, color=LIGHT_NAVY, align=PP_ALIGN.RIGHT, margin=0)


BUILDERS: dict[int, Callable] = {
    1: _cover,
    2: _pain,
    3: _workflow,
    4: _stage,
    5: _screenshots,
    6: _architecture,
    7: _comparison,
    8: _compliance,
    9: _pilot,
    10: _market,
    11: _pricing,
    12: _workflow,
    13: _founder,
    14: _timeline,
    15: _finance,
    16: _cta,
}


def build_presentation(output: str | Path, assets_dir: str | Path, *, contact=None) -> Path:
    """Build and save the complete editable BP V6 deck."""

    assets = Path(assets_dir)
    validate_assets(assets)
    destination = Path(output)
    destination.parent.mkdir(parents=True, exist_ok=True)
    prs = new_presentation()
    contact_data = dict(contact or {})
    for spec in SLIDES:
        BUILDERS[spec.number](prs, spec, assets, contact_data)
    prs.core_properties.title = "逸飞AI智眼系统｜创业大赛商业计划书 V6"
    prs.core_properties.subject = "AI视频监管闭环平台"
    prs.core_properties.author = COMPANY["name"]
    prs.core_properties.comments = "数字与表述由 tools/bp_v6/content.py 统一约束。"
    prs.save(destination)
    return destination


def read_contact(workbook_path: str | Path) -> dict[str, str]:
    """Read the existing application contact fields without changing the workbook."""

    workbook = load_workbook(workbook_path, read_only=True, data_only=True)
    try:
        sheet = workbook.worksheets[0]
        values = (sheet["P3"].value, sheet["Q3"].value, sheet["R3"].value)
    finally:
        workbook.close()
    keys = ("name", "phone", "email")
    return {key: str(value).strip() for key, value in zip(keys, values, strict=True) if value}


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--assets", type=Path, required=True)
    parser.add_argument("--workbook", type=Path)
    args = parser.parse_args()
    contact = read_contact(args.workbook) if args.workbook else {}
    result = build_presentation(args.output, args.assets, contact=contact)
    print(result)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
