from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools.bp_v6.assets import EXPECTED_ASSETS
from tools.bp_v6.build_pptx import build_presentation
from tools.bp_v6.content import SLIDES
from tools.bp_v6.theme import SLIDE_HEIGHT, SLIDE_WIDTH


def _asset_set(directory: Path) -> None:
    directory.mkdir()
    for index, name in enumerate(EXPECTED_ASSETS):
        Image.new("RGB", (1600, 900), (12 + index * 30, 32, 52)).save(
            directory / name, format="PNG"
        )


def _all_text(prs: Presentation) -> str:
    return "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def test_builds_sixteen_page_editable_deck(tmp_path):
    assets = tmp_path / "assets"
    _asset_set(assets)
    output = tmp_path / "bp-v6.pptx"

    build_presentation(
        output,
        assets,
        contact={"name": "项目联系人", "phone": "13800000000", "email": "bp@example.com"},
    )

    prs = Presentation(output)
    assert len(prs.slides) == 16
    assert prs.slide_width == SLIDE_WIDTH
    assert prs.slide_height == SLIDE_HEIGHT
    all_text = _all_text(prs)
    for spec in SLIDES:
        assert spec.title in all_text
    assert "试点验收目标，不是历史业绩" in all_text
    assert "公司成立以来收入\n0万元" in all_text
    assert "融资需求\n500万元" in all_text
    assert "P95≤10秒" in all_text
    assert "确认前不会执行任何动作" in all_text
    assert sum(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        for shape in prs.slides[4].shapes
    ) == 3


def test_deck_avoids_unapproved_claims(tmp_path):
    assets = tmp_path / "assets"
    _asset_set(assets)
    output = tmp_path / "bp-v6.pptx"
    build_presentation(output, assets, contact={})
    text = _all_text(Presentation(output))

    for forbidden in (
        "100%自研",
        "完全自主知识产权",
        "国内领先",
        "已通过等保",
        "已服务多家单位",
        "已签约",
        "已落地",
    ):
        assert forbidden not in text
